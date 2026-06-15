"""產生「熱磁材料實驗驗證計畫」簡報（~27 頁，繁中）。輸出 docs/實驗驗證計畫_簡報.pptx"""
from pathlib import Path
import lxml.etree as etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ASSETS = Path("docs/plan_assets")
OUT = Path("docs/實驗驗證計畫_簡報.pptx")
SW, SH = Inches(13.333), Inches(7.5)
NAVY=RGBColor(0x1F,0x3A,0x5F); TEAL=RGBColor(0x2A,0x9D,0x8F); ORANGE=RGBColor(0xE7,0x6F,0x51)
GOLD=RGBColor(0xE9,0xC4,0x66); DARK=RGBColor(0x26,0x2B,0x33); GREY=RGBColor(0x5A,0x63,0x70)
LIGHT=RGBColor(0xF2,0xF4,0xF7); WHITE=RGBColor(0xFF,0xFF,0xFF); PURPLE=RGBColor(0x6A,0x4C,0x93)
FONT="WenQuanYi Zen Hei"
prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH; BLANK=prs.slide_layouts[6]

def _sf(run,size,bold=False,color=DARK):
    run.font.name=FONT; run.font.size=Pt(size); run.font.bold=bold; run.font.color.rgb=color
    rPr=run._r.get_or_add_rPr()
    ea=rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None: ea=etree.SubElement(rPr,"{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface",FONT)

def rect(s,x,y,w,h,c):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=c
    sp.line.fill.background(); sp.shadow.inherit=False; return sp

def tb(s,x,y,w,h,lines,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    t=s.shapes.add_textbox(x,y,w,h).text_frame; t.word_wrap=True; t.vertical_anchor=anchor
    for i,(txt,sz,b,c,sa) in enumerate(lines):
        p=t.paragraphs[0] if i==0 else t.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.line_spacing=1.1
        r=p.add_run(); r.text=txt; _sf(r,sz,b,c)
    return t

def header(s,kick,title,n):
    rect(s,0,0,Inches(0.18),SH,TEAL)
    tb(s,Inches(0.55),Inches(0.28),Inches(11.5),Inches(0.4),[(kick,13,True,TEAL,0)])
    tb(s,Inches(0.55),Inches(0.58),Inches(12.0),Inches(0.9),[(title,28,True,NAVY,0)])
    rect(s,Inches(0.55),Inches(1.46),Inches(2.0),Pt(3),GOLD)
    tb(s,Inches(12.4),Inches(0.28),Inches(0.7),Inches(0.4),[(f"{n:02d}",13,True,GREY,0)],align=PP_ALIGN.RIGHT)

def bullets(s,x,y,w,h,items,size=17,gap=9):
    t=s.shapes.add_textbox(x,y,w,h).text_frame; t.word_wrap=True
    for i,it in enumerate(items):
        lvl=it[0]; txt=it[1]; col=it[2] if len(it)>2 else DARK
        p=t.paragraphs[0] if i==0 else t.add_paragraph()
        p.space_after=Pt(gap); p.line_spacing=1.1
        r=p.add_run(); r.text=("●  " if lvl==0 else "–  ")+txt
        _sf(r,size if lvl==0 else size-2,lvl==0,col)
    return t

def img(s,p,x,y,w=None,h=None):
    return s.shapes.add_picture(str(p),x,y,width=w) if w else s.shapes.add_picture(str(p),x,y,height=h)

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    if bg!=WHITE: rect(s,0,0,SW,SH,bg)
    return s

def table(s,x,y,col_w,rows,fs=12,rh=0.46,hdr=NAVY):
    yy=y
    for ri,row in enumerate(rows):
        xx=x; h=Inches(rh)
        for ci,cell in enumerate(row):
            w=Inches(col_w[ci])
            bg=hdr if ri==0 else (LIGHT if ri%2 else WHITE)
            rect(s,xx,yy,w,h,bg)
            tb(s,xx+Inches(0.05),yy,w-Inches(0.1),h,
               [(cell,fs,ri==0,WHITE if ri==0 else DARK,0)],anchor=MSO_ANCHOR.MIDDLE)
            xx+=w
        yy+=h

# 01 封面
s=slide(NAVY)
rect(s,0,Inches(5.2),SW,Inches(0.12),GOLD)
tb(s,Inches(0.9),Inches(1.4),Inches(11.5),Inches(0.5),[("EXPERIMENTAL VALIDATION CAMPAIGN",15,True,GOLD,0)])
tb(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(2.0),
   [("熱磁材料實驗驗證計畫",44,True,WHITE,6),("小規模先導 → 階段閘控 → 規模化擴張",28,True,LIGHT,0)])
tb(s,Inches(0.9),Inches(4.1),Inches(11.5),Inches(1.0),
   [("把模型的『相對可信』收斂成『絕對可信』；先以最小可行先導去風險，避免重蹈覆轍",17,False,LIGHT,0)])
tb(s,Inches(0.9),Inches(6.3),Inches(11.5),Inches(0.5),[("alloy_engine · 實驗計畫 + 詳細需求 · 2026",13,False,GOLD,0)])

# 02 一句話
s=slide(); header(s,"EXECUTIVE SUMMARY","計畫一句話",2)
tb(s,Inches(0.9),Inches(2.0),Inches(11.6),Inches(1.6),
   [("模型已把第一性原理 + 公開資料能做的都做完，剩下唯有實測能解。",22,False,DARK,6),
    ("策略：先用最易取得的『純 Gd』打通整條量測—回填流程（最小成本驗證可行性），",22,True,NAVY,6),
    ("通過決策門後再逐步擴張到一階材料、GA 候選、複合與整機原型。",22,False,DARK,0)])
img(s,ASSETS/"funnel.png",Inches(2.6),Inches(3.9),w=Inches(8.1))

# 03 為什麼做實驗
s=slide(); header(s,"WHY","為什麼一定要做實驗",3)
bullets(s,Inches(0.9),Inches(1.85),Inches(11.6),Inches(5),
   [(0,"模型現況：相對可信、絕對待測。"),
    (1,"真實 Tc baseline R²=0.78（已用 NEMAD 真實資料）。"),
    (1,"Br 溫度修正後 bias 僅 −0.12T（多為 0K-vs-室溫差，非系統性低估）。"),
    (1,"發電側功率密度是理想化上界（比真實原型高 ~10×）。"),
    (1,"複合 connectivity 影響增益 ~43%，但定性結論穩健。"),
    (0,"這些『絕對值待收斂』的缺口，只有真實樣品的 M-H/DSC/SEM 能補。",ORANGE),
    (0,"目標：用最小規模、可控成本，把絕對值校準，讓設計能真正落地。",NAVY)])

# 04 核心理念
s=slide(); header(s,"PRINCIPLE","核心理念：先導 + 決策門 + 擴張",4)
img(s,ASSETS/"derisk.png",Inches(0.7),Inches(1.9),w=Inches(12.0))
bullets(s,Inches(0.9),Inches(5.5),Inches(11.6),Inches(1.6),
   [(0,"每階段先小規模驗證『流程＋模型量級』，過了決策門才放大投入。"),
    (0,"任一門未過 → 回上一階段修模型/製程，不貿然擴張——這就是避免重蹈覆轍。",ORANGE)],size=15,gap=7)

# 05 總體目標
s=slide(); header(s,"OBJECTIVES","總體目標（SMART）",5)
table(s,Inches(0.7),Inches(1.9),[2.0,7.5,2.8],
   [["目標","內容","量化判準"],
    ["G-A 流程可信","打通製樣→量測→回填全鏈","Gd ΔM/ΔS 與文獻 ±20% 內"],
    ["G-B 模型校準","用實測收斂 ΔM/w/connectivity","關鍵參數誤差 < 25%"],
    ["G-C 端到端","GA 候選的預測 vs 實測","Tc 誤差 < 50°C、排序正確"],
    ["G-D 絕對校準","整機原型錨定發電側","P/V、η 量級對齊文獻原型"]],fs=13)
tb(s,Inches(0.7),Inches(5.6),Inches(12),Inches(0.8),
   [("每個目標都對應一個可量測的『過/不過』判準，避免主觀。",15,True,GREY,0)])

# 06 成功判準 / KPI
s=slide(); header(s,"KPI","成功判準與 KPI",6)
bullets(s,Inches(0.9),Inches(1.9),Inches(11.6),Inches(5),
   [(0,"先導 KPI（Phase 0，最關鍵）：純 Gd 的 ΔM(T)、ΔS_M、Tc 與文獻吻合（±20%）。"),
    (1,"若連最單純的二階材料都對不上 → 是量測/流程問題，必須先解決再前進。"),
    (0,"模型校準 KPI：回填後，reference_materials 與 design_tmg 的相對排序不變、絕對值收斂。"),
    (0,"端到端 KPI：GA 推薦合金實測 Tc 誤差 < 50°C；Top-N 排序與實測一致。"),
    (0,"整機 KPI：原型 P/V、η 與文獻 TMG 原型同量級（收斂 D12 的 ~10× 落差）。"),
    (0,"全程 KPI：每階段資料入庫、可回填、可復現（對應 git 內腳本）。",TEAL)])

# 07 驗證對象
s=slide(); header(s,"WHAT TO VALIDATE","要驗證什麼（模型 → 真實）",7)
table(s,Inches(0.7),Inches(1.9),[3.2,4.0,5.0],
   [["模型量","現況（假設/估計）","實測校準"],
    ["ΔM(工作溫度)","平均場 m(T/Tc) 推估","VSM M-H 迴線"],
    ["一階銳度 w (D5)","假設值","M(T) 斜率擬合"],
    ["複合 connectivity (D6)","假設 0.7（敏感 43%）","SEM 相連通度"],
    ["κ / Cp / ΔS_M","純元素混合估計","雷射閃光 / DSC"],
    ["氫化 Tc 上修 (D8)","heuristic 模型","氫化前後 VSM"],
    ["整機 W/η/P (D12)","理想化上界","原型實測"]],fs=12.5,rh=0.52)

# 08 樣品計畫
s=slide(); header(s,"SAMPLES","樣品計畫（最小可驗證集）",8)
table(s,Inches(0.6),Inches(1.85),[1.0,2.6,3.0,2.0,3.2],
   [["#","材料","目的","Tc","取得"],
    ["M1","純 Gd","二階基準、流程打通","~293K","商購（高純）"],
    ["M2","La(Fe,Si)13(H)","一階 + 氫化驗證(D8)","200→340K","委外/合作合成"],
    ["M3","(Mn,Fe)2(P,Si)","無稀土一階、ΔS 最大","~280K","委外/合作合成"],
    ["M4","GA 候選 Fe 系 1–2","端到端模型驗證","依搜尋","電弧熔煉自製"]],fs=12.5,rh=0.55)
tb(s,Inches(0.6),Inches(5.6),Inches(12),Inches(1.0),
   [("選材邏輯：覆蓋二階/一階、有/無稀土、文獻基準/模型自選——以最少樣品覆蓋最多物理區。",15,True,GREY,0)])

# 09 樣品取得策略
s=slide(); header(s,"ACQUISITION","樣品取得策略：買 vs 自製 vs 委外",9)
table(s,Inches(0.7),Inches(1.9),[2.2,3.6,3.4,2.3],
   [["方式","適用","優點","風險/成本"],
    ["商購","純 Gd、純元素","快、純度保證、低風險","成分受限"],
    ["自製(電弧熔煉)","Fe 系 GA 候選","成分自由、快迭代","需熔煉+退火設備"],
    ["委外/學界合作","La-Fe-Si、Mn-Fe-P","拿到難製一階相","週期長、需協調"]],fs=12.5,rh=0.55)
bullets(s,Inches(0.7),Inches(5.0),Inches(11.8),Inches(1.8),
   [(0,"先導 Phase 0 只需『商購純 Gd』→ 零合成風險、最快啟動。",TEAL),
    (0,"一階材料優先找有 La-Fe-Si/Mn-Fe-P 經驗的實驗室合作，避免自行踩相形成的坑。",ORANGE)],size=15,gap=7)

# 10 Phase 0
s=slide(); header(s,"PHASE 0","先導實驗：純 Gd（最小可行）",10)
bullets(s,Inches(0.9),Inches(1.9),Inches(6.0),Inches(5),
   [(0,"為何選 Gd：商購可得、二階、文獻 M-H/ΔS_M 完整。"),
    (0,"範圍：1 個樣品、2 項核心量測（VSM + DSC）。"),
    (0,"目的：打通『製樣→量測→回填→比對文獻』整條鏈。"),
    (0,"成本/時間：最低（~2 個月、1 樣）。"),
    (0,"產出：驗證我們的量測與回填流程本身可信。",TEAL)],size=16)
rect(s,Inches(7.2),Inches(2.0),Inches(5.5),Inches(3.6),LIGHT)
tb(s,Inches(7.4),Inches(2.2),Inches(5.1),Inches(3.3),
   [("Phase 0 為什麼關鍵",17,True,NAVY,10),
    ("先用『答案已知』的 Gd 校準儀器與流程，",14,False,DARK,8),
    ("確認我們量得準、回填對。",14,False,DARK,8),
    ("若連 Gd 都對不上 → 先修流程，",14,True,ORANGE,8),
    ("不浪費昂貴的一階樣品。",14,True,ORANGE,0)])

# 11 Phase 0 步驟
s=slide(); header(s,"PHASE 0 STEPS","先導步驟與量測",11)
table(s,Inches(0.7),Inches(1.9),[1.0,3.4,3.6,3.0],
   [["步","活動","量測/輸出","對應模型"],
    ["1","採購高純 Gd + 切樣","EDS/XRD 確認純度相","—"],
    ["2","VSM M-H @ 多溫度","ΔJ(T)、磁滯面積","ΔM / 磁滯損耗"],
    ["3","DSC 過 Tc","ΔS_M、Cp、Tc","reference_materials"],
    ["4","M(T) 擬合","二階曲線形狀","驗證平均場 m(T/Tc)"],
    ["5","比對文獻 + 回填","±20% 判定","決策門 G0"]],fs=12.5,rh=0.5)

# 12 決策門 G0
s=slide(); header(s,"GATE G0","決策門 G0：先導通過了嗎？",12)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.6),Inches(3),
   [(0,"通過條件：Gd 的 ΔM(T)、ΔS_M、Tc 與文獻吻合（±20%），且回填流程可復現。",TEAL),
    (0,"通過 → 進 Phase 1（一階材料），開始投入較貴的合成樣品。"),
    (0,"未過 → 診斷是量測（校準/退火/氧化）還是模型（公式/單位）問題，修正後重做先導。",ORANGE),
    (0,"原則：用最便宜的一步，攔下最貴的錯誤——避免重蹈覆轍。",NAVY)])

# 13 Phase 1
s=slide(); header(s,"PHASE 1","一階材料：La-Fe-Si / Mn-Fe-P",13)
bullets(s,Inches(0.9),Inches(1.9),Inches(11.6),Inches(4.6),
   [(0,"對象：La(Fe,Si)13 與 (Mn,Fe)2(P,Si)——真正高 ΔS_M 的一階 MCE 主力。"),
    (0,"重點量測：M(T) 銳度 → 校 D5 logistic 寬度 w；ΔS_M → 校 reference_materials。"),
    (0,"挑戰：一階相需正確退火/相形成；脆、需小心製樣（與 D9 脆性懲罰呼應）。"),
    (1,"故優先委外/與有經驗的實驗室合作，降低製程風險。"),
    (0,"規模：2–3 個樣品；先各 1 個成分點驗證，再視結果加成分掃描。",TEAL)])

# 14 Phase 1 氫化
s=slide(); header(s,"PHASE 1 · D8","氫化 Tc 上修驗證",14)
bullets(s,Inches(0.9),Inches(1.9),Inches(11.6),Inches(4.6),
   [(0,"模型：hydrogenation_tc_shift_K() 預測 La-Fe-Si 吸氫後 Tc 上修 ~50–150K。"),
    (0,"實驗：同一 La-Fe-Si 樣品，氫化前後各量 VSM/DSC，比較 Tc 位移。"),
    (0,"判準：氫化後 Tc 上移方向正確、量級吻合（±30%）。"),
    (0,"價值：驗證 D8 把『近室溫材料』調到『廢熱溫區』的可行性（解 D10 溫區匹配）。",NAVY)])

# 15 Phase 2
s=slide(); header(s,"PHASE 2","GA 候選端到端 + 複合 connectivity",15)
bullets(s,Inches(0.9),Inches(1.9),Inches(11.6),Inches(4.8),
   [(0,"GA 候選：用烘焙好的 bundle_real_tc（真實 Tc）跑 run_search，取 Top 1–2 Fe 系配方。"),
    (1,"電弧熔煉自製 → 量 Tc/Br/σy → 比對模型預測（端到端 sim-to-real 驗證）。"),
    (0,"複合 connectivity（D6 最敏感參數）：製備高κ基底+一階相複合，SEM 量真實連通度。"),
    (1,"回填 composite.py connectivity，收斂功率增益的絕對值（現為 ×7.6–×11.8 帶）。"),
    (0,"規模：~6 樣（GA 候選 + 複合系列）。",TEAL)])

# 16 Phase 3
s=slide(); header(s,"PHASE 3","整機 TMG 原型（發電側絕對校準）",16)
bullets(s,Inches(0.9),Inches(1.9),Inches(11.6),Inches(4.8),
   [(0,"目標：自製小型 TMG 原型，量測真實 W / η / P-V。"),
    (0,"用前面校準好的最佳材料（一階 + 複合），在實際磁路/熱循環下運轉。"),
    (0,"把結果新增為 reference_devices.py 的『本工作』錨點 → 收斂 D12 的 ~10× 落差。"),
    (0,"這是最大工程量、放最後做——前三階段全綠燈才啟動。",ORANGE),
    (0,"成果：整機絕對功率/效率可預測，模型從『相對工具』升級為『絕對預測』。",NAVY)])

# 17 量測矩陣
s=slide(); header(s,"MEASUREMENT MATRIX","量測矩陣（設備 × 測項）",17)
table(s,Inches(0.6),Inches(1.85),[3.0,2.4,3.0,3.6],
   [["設備","測項","輸出","可得性"],
    ["EDS / XRD","成分 / 相","純度、相確認","高（常見）"],
    ["VSM / PPMS","M-H 迴線 / M(T)","ΔJ、磁滯、w","中（需磁量測）"],
    ["DSC","ΔS_M / Cp / Tc","熱量、相變峰","中"],
    ["雷射閃光","熱擴散 α","κ","中"],
    ["SEM","微結構 / 連通度","connectivity、φ","高"],
    ["自製 TMG 台","W / η / P","整機效能","低（工程）"]],fs=12.5,rh=0.5)

# 18 回填對應
s=slide(); header(s,"CLOSED LOOP","量到什麼 → 回填模型哪裡",18)
img(s,ASSETS/"backfill.png",Inches(1.4),Inches(1.85),w=Inches(10.5))

# 19 時程
s=slide(); header(s,"TIMELINE","時程（階段化 Gantt）",19)
img(s,ASSETS/"gantt.png",Inches(1.0),Inches(1.95),w=Inches(11.4))

# 20 里程碑
s=slide(); header(s,"MILESTONES","里程碑與決策門",20)
table(s,Inches(0.7),Inches(1.9),[1.4,2.2,5.0,3.0],
   [["月","里程碑","內容","決策門"],
    ["M2","先導完成","Gd 流程打通、回填驗證","G0 go/no-go"],
    ["M5","一階驗證","La-Fe-Si/Mn-Fe-P + 氫化","—"],
    ["M8","端到端","GA 候選 + 複合 connectivity","G1 go/no-go"],
    ["M12","整機原型","發電側絕對校準","計畫收尾/續期"]],fs=13,rh=0.55)
tb(s,Inches(0.7),Inches(5.5),Inches(12),Inches(0.8),
   [("兩個 go/no-go 門（G0/G1）是去風險核心：沒過不放大投入。",15,True,ORANGE,0)])

# 21 設備需求
s=slide(); header(s,"RESOURCES","設備需求與取得管道",21)
bullets(s,Inches(0.9),Inches(1.9),Inches(11.6),Inches(4.8),
   [(0,"自有/校內共儀：XRD、SEM、DSC、VSM/PPMS——多數材料系所具備，先盤點。"),
    (0,"合成：電弧熔煉爐 + 退火爐（Fe 系自製）；一階相委外或合作。"),
    (0,"氫化：H2 氣氛爐（La-Fe-Si-H），需安全規範——可委外。"),
    (0,"整機台：自製（磁路 + 換熱 + 線圈 + 量測），工程量大，Phase 3 才投入。"),
    (0,"策略：先用共儀與委外把固定成本壓到最低，驗證可行再考慮自建。",TEAL)])

# 22 預算
s=slide(); header(s,"BUDGET","預算 / 資源粗估（相對量級）",22)
table(s,Inches(0.7),Inches(1.9),[2.6,3.0,3.2,2.5],
   [["階段","樣品/合成","量測（共儀/委外）","相對量級"],
    ["Phase 0","商購 Gd（低）","VSM+DSC（低）","★ 最低"],
    ["Phase 1","委外一階（中）","VSM+DSC+氫化（中）","★★"],
    ["Phase 2","自製+複合（中）","+SEM+κ（中）","★★★"],
    ["Phase 3","原型製作（高）","整機量測（高）","★★★★"]],fs=13,rh=0.55)
tb(s,Inches(0.7),Inches(5.6),Inches(12),Inches(0.8),
   [("精確數字依在地共儀/委外報價而定；本表為相對量級，供排序與分期申請經費。",15,True,GREY,0)])

# 23 風險
s=slide(); header(s,"RISK","風險登錄與緩解（避免重蹈覆轍）",23)
table(s,Inches(0.6),Inches(1.85),[3.6,4.2,3.8],
   [["風險","影響","緩解"],
    ["一階相未正確形成","量到的不是目標相","委外/合作 + XRD 先確認相"],
    ["稀土氧化/氫脆","樣品劣化、數據失真","惰性氣氛、鍍層、儘速量測"],
    ["量測未校準","絕對值不可信","Phase 0 用 Gd 先校準儀器"],
    ["範圍蔓延","成本失控","決策門閘控、分期投入"],
    ["GA 候選不可製造","白做樣品","先過 D9 脆性/可製造性篩選"]],fs=12,rh=0.55)

# 24 規模化
s=slide(); header(s,"SCALE-UP","規模化路徑：先小後大",24)
img(s,ASSETS/"scale_ramp.png",Inches(1.4),Inches(1.9),w=Inches(8.2))
bullets(s,Inches(9.7),Inches(2.1),Inches(3.4),Inches(4.5),
   [(0,"Phase 0：1 樣驗流程。"),
    (0,"Phase 1：3 樣一階。"),
    (0,"Phase 2：6 樣 GA+複合。"),
    (0,"Phase 3：原型。"),
    (0,"投入隨『已驗證的把握』遞增。",TEAL)],size=14,gap=9)

# 25 團隊
s=slide(); header(s,"TEAM","團隊與角色（最小編制）",25)
table(s,Inches(0.7),Inches(1.9),[3.0,4.5,3.5],
   [["角色","職責","階段"],
    ["材料製備","合成/退火/委外協調","全程"],
    ["磁/熱量測","VSM/DSC/κ + 數據","全程"],
    ["建模回填","實測→模型校準→復現","全程"],
    ["整機工程","原型設計與量測","Phase 3"]],fs=13,rh=0.55)
tb(s,Inches(0.7),Inches(5.4),Inches(12),Inches(0.9),
   [("可由 2–3 人 + 共儀起步；建模回填角色與既有 alloy_engine 腳本直接銜接。",15,True,GREY,0)])

# 26 交付物
s=slide(); header(s,"DELIVERABLES","交付物與下一步",26)
bullets(s,Inches(0.9),Inches(1.9),Inches(11.6),Inches(4.8),
   [(0,"每階段：量測原始數據 + 回填後的模型參數 + 復現腳本（入 git）。"),
    (0,"先導報告（Phase 0）：Gd 流程驗證 + G0 決策建議。"),
    (0,"校準後模型：reference_materials / composite / device 的實測版。"),
    (0,"最終：整機原型對標報告，收斂發電側絕對校準。"),
    (0,"立即下一步：① 盤點可用共儀 ② 採購純 Gd ③ 跑 run_search 產 GA 候選清單。",NAVY)])

# 27 結語
s=slide(NAVY)
rect(s,0,Inches(3.4),SW,Inches(0.1),GOLD)
tb(s,Inches(0.9),Inches(1.6),Inches(11.5),Inches(1.6),
   [("先導去風險，分期擴張",38,True,WHITE,8),("用最小成本，把模型校成可信的設計工具",24,False,LIGHT,0)])
t=s.shapes.add_textbox(Inches(0.9),Inches(4.0),Inches(11.5),Inches(2.6)).text_frame; t.word_wrap=True
for i,x in enumerate([
   "Phase 0 純 Gd 打通流程 → G0 → 一階材料 → GA 候選/複合 → G1 → 整機原型。",
   "每一步都有可量測的決策門，沒過不放大投入——避免重蹈覆轍。",
   "量到什麼都對應回填模型哪裡，形成閉環校準。",
   "詳細需求、規格、時程、預算見隨附 Word 文件。"]):
    p=t.paragraphs[0] if i==0 else t.add_paragraph(); p.space_after=Pt(11)
    r=p.add_run(); r.text="●  "+x; _sf(r,16,False,LIGHT)

prs.save(str(OUT))
print("saved",OUT,"slides",len(prs.slides._sldIdLst))
