"""初學者向簡報：用最白話、從基礎數學，快速看懂整個專案。輸出 docs/熱磁專案_白話入門.pptx"""
from pathlib import Path
import lxml.etree as etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BEG = Path("docs/beginner_assets")
OUT = Path("docs/熱磁專案_白話入門.pptx")
SW, SH = Inches(13.333), Inches(7.5)
NAVY=RGBColor(0x1F,0x3A,0x5F); TEAL=RGBColor(0x2A,0x9D,0x8F); ORANGE=RGBColor(0xE7,0x6F,0x51)
GOLD=RGBColor(0xE9,0xC4,0x66); DARK=RGBColor(0x26,0x2B,0x33); GREY=RGBColor(0x5A,0x63,0x70)
LIGHT=RGBColor(0xF2,0xF4,0xF7); WHITE=RGBColor(0xFF,0xFF,0xFF); RED=RGBColor(0xC1,0x12,0x1F)
FONT="WenQuanYi Zen Hei"
prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH; BLANK=prs.slide_layouts[6]

def _sf(run,size,bold=False,color=DARK):
    run.font.name=FONT; run.font.size=Pt(size); run.font.bold=bold; run.font.color.rgb=color
    rPr=run._r.get_or_add_rPr()
    ea=rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None: ea=etree.SubElement(rPr,"{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface",FONT)

def rect(s,x,y,w,h,c):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=c
    sp.line.fill.background(); sp.shadow.inherit=False; return sp

def tb(s,x,y,w,h,lines,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    t=s.shapes.add_textbox(x,y,w,h).text_frame; t.word_wrap=True; t.vertical_anchor=anchor
    for i,(txt,sz,b,c,sa) in enumerate(lines):
        p=t.paragraphs[0] if i==0 else t.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.line_spacing=1.15
        r=p.add_run(); r.text=txt; _sf(r,sz,b,c)
    return t

def header(s,kick,title,n):
    rect(s,0,0,Inches(0.18),SH,TEAL)
    tb(s,Inches(0.55),Inches(0.30),Inches(11.5),Inches(0.4),[(kick,14,True,TEAL,0)])
    tb(s,Inches(0.55),Inches(0.64),Inches(12.1),Inches(0.95),[(title,29,True,NAVY,0)])
    rect(s,Inches(0.55),Inches(1.56),Inches(2.0),Pt(3),GOLD)
    tb(s,Inches(12.4),Inches(0.30),Inches(0.7),Inches(0.4),[(f"{n:02d}",13,True,GREY,0)],align=PP_ALIGN.RIGHT)

def bullets(s,x,y,w,h,items,size=18,gap=12):
    t=s.shapes.add_textbox(x,y,w,h).text_frame; t.word_wrap=True
    for i,it in enumerate(items):
        lvl=it[0]; txt=it[1]; col=it[2] if len(it)>2 else DARK
        p=t.paragraphs[0] if i==0 else t.add_paragraph(); p.space_after=Pt(gap); p.line_spacing=1.15
        r=p.add_run(); r.text=("●  " if lvl==0 else "    – ")+txt
        _sf(r,size if lvl==0 else size-3,lvl==0,col)
    return t

def img(s,p,x,y,w): return s.shapes.add_picture(str(p),x,y,width=w)
def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    if bg!=WHITE: rect(s,0,0,SW,SH,bg)
    return s

def bigmath(s,formula,plain_lines,y=Inches(2.4)):
    """大字公式 + 白話逐項解釋。"""
    rect(s,Inches(1.4),y,Inches(10.5),Inches(1.15),LIGHT)
    tb(s,Inches(1.4),y,Inches(10.5),Inches(1.15),[(formula,30,True,NAVY,0)],
       align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    bullets(s,Inches(1.6),y+Inches(1.5),Inches(10),Inches(3),plain_lines,size=17,gap=10)

# 01 封面
s=slide(NAVY)
rect(s,0,Inches(5.3),SW,Inches(0.12),GOLD)
tb(s,Inches(0.9),Inches(1.7),Inches(11.5),Inches(2.2),
   [("10 分鐘白話看懂", 30, True, GOLD, 6),
    ("熱磁發電 + AI 找材料", 46, True, WHITE, 8),
    ("從最基礎的數學講起，不用背公式也能懂", 22, False, LIGHT, 0)])
tb(s,Inches(0.9),Inches(6.4),Inches(11.5),Inches(0.5),[("給完全新手的入門簡報 · alloy_engine", 14, False, GOLD, 0)])

# 02 在幹嘛
s=slide(); header(s,"這專案在幹嘛？","一句話：把『沒人要的熱』變成電",2)
img(s,BEG/"flow.png",Inches(1.6),Inches(2.0),Inches(10.2))
bullets(s,Inches(1.0),Inches(5.2),Inches(11.5),Inches(1.8),
   [(0,"工廠、引擎到處有『廢熱』排掉很可惜——能不能撿來發電？"),
    (0,"用一種特別的材料：它一被加熱，磁性就會變；磁性一變，就能生電。",TEAL)],size=16)

# 03 磁鐵是什麼
s=slide(); header(s,"先備知識 1","磁鐵到底是什麼？",3)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.5),
   [(0,"想像材料裡每個原子都是一根『小磁針』。"),
    (0,"這些小磁針方向一致 → 整體就有磁性（吸鐵）。"),
    (0,"方向亂七八糟 → 互相抵消 → 沒磁性。"),
    (0,"重點：『溫度』會影響這些小磁針——越熱，越容易被晃亂。",NAVY)])

# 04 居禮溫度
s=slide(); header(s,"先備知識 2","居禮溫度 Tc = 磁性的開關",4)
img(s,BEG/"tc_switch.png",Inches(1.5),Inches(1.95),Inches(8.0))
bullets(s,Inches(9.7),Inches(2.3),Inches(3.4),Inches(4),
   [(0,"加熱超過某個溫度，"),(0,"小磁針被晃亂 →"),(0,"磁性『消失』。",RED),
    (0,"這個溫度就叫"),(0,"居禮溫度 Tc。",NAVY)],size=15,gap=9)

# 05 核心點子
s=slide(); header(s,"核心點子","為什麼『磁性會變』就能發電？",5)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.5),
   [(0,"200 年前法拉第就發現：磁場一變化，旁邊的線圈就會生出電流。"),
    (0,"所以只要讓材料『一下有磁、一下沒磁』反覆切換——線圈就一直發電。"),
    (0,"怎麼切換？把材料在『冷』和『熱』之間來回（冷端有磁、熱端沒磁）。",TEAL),
    (0,"這就是熱磁發電：用溫度差，驅動磁性反覆變化 → 發電。",NAVY)])

# 06 基礎數學①
s=slide(); header(s,"基礎數學 1（別怕）","發多少電？看『磁性變化』有多大",6)
bigmath(s,"W  =  μ₀ ∮ H dM",
   [(0,"W = 一次循環發出的『功』（能量）——就是我們要的電。"),
    (0,"∮ H dM = 磁性變化（M）× 外加磁場推力（H）累積起來。",TEAL),
    (0,"白話：磁性變化越大、磁場推力越強 → 發的電越多。",NAVY),
    (0,"μ₀ 只是個固定常數，不用管它。",GREY)])

# 07 基礎數學②
s=slide(); header(s,"基礎數學 2","效率？= 拿到的電 ÷ 投入的熱",7)
bigmath(s,"η  =  W ÷ Q",
   [(0,"η（讀 eta）= 效率，就像汽車的『油耗』。"),
    (0,"W = 發出的電（拿到的）；Q = 投入的熱（花掉的）。",TEAL),
    (0,"白話：花一堆熱，只發一點電 → 效率低；反之效率高。"),
    (0,"這專案發現：大部分熱其實『白白浪費』掉了——下一頁解釋。",ORANGE)])

# 08 卡諾天花板
s=slide(); header(s,"先講清楚","再厲害也有『天花板』",8)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.5),
   [(0,"物理鐵律（卡諾）：任何熱機的效率，都有一個上限。"),
    (0,"上限由『冷熱溫差』決定——溫差越大，上限越高。"),
    (0,"廢熱通常溫差不大（如 150°C vs 室溫）→ 天花板本來就不高。",ORANGE),
    (0,"所以我們不是要違反物理，而是要『盡量逼近』這個天花板。",NAVY)])

# 09 第一個難題
s=slide(); header(s,"第一個難題","大部分熱被『白白加熱』掉了",9)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.6),
   [(0,"投入的熱分兩種用途："),
    (1,"一部分真的拿去『改變磁性』→ 有用，能發電。",TEAL),
    (1,"一大部分只是『把材料整個加熱』→ 白白浪費（叫顯熱）。",ORANGE),
    (0,"廢熱發電裡，浪費的那部分常佔 80–95%！效率因此超低。",RED),
    (0,"救星：『回熱器』——把材料降溫時放出的熱，回收去預熱下一批。",NAVY)])

# 10 為什麼要電腦
s=slide(); header(s,"換個角度","材料配方太多，人試不完",10)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.5),
   [(0,"好材料要同時滿足：對的居禮溫度、磁性強、導熱好、做得出來…"),
    (0,"我們有 14 種元素可調比例 → 配方組合是天文數字。"),
    (0,"一個一個做實驗試？幾輩子也試不完。",ORANGE),
    (0,"所以：先用電腦（AI）大量篩選，挑出最有希望的幾個再做實驗。",NAVY)])

# 11 AI 預測
s=slide(); header(s,"AI 第一招","用『代理模型』猜性質",11)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.5),
   [(0,"像房價預測：給『坪數、地點、屋齡』→ 猜『房價』。"),
    (0,"這裡：給『合金成分』→ 猜『居禮溫度、磁性…』。",TEAL),
    (0,"AI 看過很多例子後，就能對沒看過的新配方快速給出估計。"),
    (0,"好處：一秒可以『猜』幾十萬個配方，不必真的去做。",NAVY)])

# 12 GA
s=slide(); header(s,"AI 第二招","用『遺傳演算法』演化出好配方",12)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.5),
   [(0,"靈感來自生物演化 / 育種："),
    (1,"先隨機生一大堆配方，用 AI 評分。"),
    (1,"留下分數高的，讓它們『交配、突變』生出下一代。"),
    (1,"一代一代下來，配方越來越好。",TEAL),
    (0,"電腦每秒能評估數百萬個——很快就逼出最佳配方。",NAVY)])

# 13 大轉折
s=slide(); header(s,"最重要的轉折","一開始 AI 其實『很不準』",13)
img(s,BEG/"sim_real.png",Inches(1.6),Inches(2.0),Inches(7.0))
bullets(s,Inches(8.9),Inches(2.3),Inches(4.0),Inches(4),
   [(0,"一開始用『電腦模擬的假資料』教 AI。"),
    (0,"拿到真實世界一比 →",ORANGE),
    (0,"比亂猜還差！(左邊紅柱)",RED),
    (0,"問題不在 AI，在『教材是假的』。",NAVY)],size=14,gap=10)

# 14 救援
s=slide(); header(s,"怎麼救？","換成『真實資料』，AI 立刻變準",14)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.6),
   [(0,"找到一個公開資料庫：一萬五千筆『真實量測』的居禮溫度。"),
    (0,"用真實資料重教 AI → 準度從『比亂猜差』變成『相當準』。",TEAL),
    (0,"磁性（Br）也如法炮製：用真實的物理計算資料 → 也變準。"),
    (0,"教訓一句話：垃圾進、垃圾出；真實資料才是關鍵。",NAVY)])

# 15 誠實①
s=slide(); header(s,"我們很誠實 1","有些性質，天生就『猜不準』",15)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.6),
   [(0,"居禮溫度、磁性 → 主要由『成分』決定 → AI 猜得準。",TEAL),
    (0,"但『矯頑力、強度』→ 主要由『怎麼做的』決定（晶粒、加工…）。",ORANGE),
    (0,"同樣成分，做法不同，結果差很多 → 光看成分本來就猜不準。"),
    (0,"我們沒有硬掰，而是老實標明：『這兩個只能當參考，不保證準』。",NAVY)])

# 16 瓶頸鏈
s=slide(); header(s,"我們學到的","卡關的地方，一個接一個",16)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.6),
   [(0,"① 顯熱：大部分熱被浪費（前面講過）。"),
    (0,"② 磁滯：磁性來回切換會『卡頓』損耗能量。"),
    (0,"③ 導熱差：熱傳不進傳不出，材料來不及冷熱切換 → 發電變慢。",ORANGE),
    (0,"每解一個，就冒出下一個——這就是真實工程的樣子。",NAVY)])

# 17 複合材料
s=slide(); header(s,"一個聰明解法","複合材料：像鋼筋混凝土",17)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.6),
   [(0,"最磁性強的材料，常常『導熱很差』（卡在瓶頸③）。"),
    (0,"點子：把『磁性強的料』 + 『導熱好的料（如銅）』混在一起。",TEAL),
    (1,"就像混凝土（耐壓）配鋼筋（耐拉）——各取所長。"),
    (0,"犧牲一點磁性，換來導熱大增 → 整體發電反而更好。",NAVY)])

# 18 誠實②
s=slide(); header(s,"我們很誠實 2","模型給的是『相對好壞』，不是保證值",18)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.6),
   [(0,"模型很會比『A 配方比 B 好』——相對排序可信。",TEAL),
    (0,"但『絕對發電量』的數字偏樂觀（理想化），實機會打折。",ORANGE),
    (0,"所以我們給的預測都『帶誤差範圍』，並標明哪些要靠實驗確認。"),
    (0,"科學的誠實：講清楚『什麼可信、什麼還不確定』。",NAVY)])

# 19 主動學習
s=slide(); header(s,"最後一步","電腦建議『先做哪個實驗』",19)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.6),
   [(0,"實驗很貴，不能亂做。電腦可以建議『先做最有價值的那幾個』。"),
    (0,"我們做了，也老實去驗證它有沒有用——"),
    (0,"結果發現：在目前資料量下，『隨機挑』反而更好！",RED),
    (0,"所以我們據實說明，沒有誇大這個功能。這也是誠實的一部分。",NAVY)])

# 20 全貌
s=slide(); header(s,"一張圖看全貌","整個專案的流水線",20)
img(s,BEG/"pipeline.png",Inches(1.1),Inches(2.2),Inches(11.2))
tb(s,Inches(0.9),Inches(5.3),Inches(11.5),Inches(1.4),
   [("配方 → AI 猜好壞 → 遺傳演算法挑最好 → 候選名單 → 做實驗驗證。",18,True,NAVY,6),
    ("AI 用『真實資料』教出來；最後一步要靠實驗室。",16,False,GREY,0)])

# 21 教會我們什麼
s=slide(); header(s,"這專案教會我們","三個帶得走的觀念",21)
bullets(s,Inches(0.9),Inches(2.0),Inches(11.5),Inches(4.6),
   [(0,"① 從『第一性原理』出發：先懂物理（磁、熱、電），再寫程式。",TEAL),
    (0,"② 『真實資料』決定成敗：假資料教出的 AI 一文不值。",NAVY),
    (0,"③ 『誠實』比好看重要：講清楚能信什麼、不能信什麼。",ORANGE)])

# 22 結語
s=slide(NAVY)
rect(s,0,Inches(3.5),SW,Inches(0.1),GOLD)
tb(s,Inches(0.9),Inches(1.7),Inches(11.5),Inches(1.6),
   [("看懂了嗎？", 40, True, WHITE, 8),("把廢熱變電，用 AI + 真實資料 + 誠實去逼近物理極限", 22, False, LIGHT, 0)])
t=prs.slides[-1].shapes.add_textbox(Inches(0.9),Inches(4.0),Inches(11.5),Inches(2.6)).text_frame; t.word_wrap=True
for i,x in enumerate([
   "熱 → 磁性變化 → 電：這是整條主線。",
   "AI 幫忙在天文數字的配方裡篩選，但教材必須是真實資料。",
   "想深入？看技術版簡報與 docs/（含論文初稿、設計案例）。"]):
    p=t.paragraphs[0] if i==0 else t.add_paragraph(); p.space_after=Pt(12)
    r=p.add_run(); r.text="●  "+x; _sf(r,17,False,LIGHT)

prs.save(str(OUT)); print("saved",OUT,"slides",len(prs.slides._sldIdLst))
