"""產生『熱磁發電機架構規格（飛輪 vs 活塞）』PPT + Word。

整合三方輸入：
  1. 網路調研 6 子系統章節 + 規格表 + 引用（/tmp/tmg_arch_synth.json，由研究工作流產出）
  2. GA 架構↔材料共同設計結果（calibrated 預設，內嵌於本檔 GA_ROWS）
  3. 引擎整機模型方程（generator_design / device_score）

輸出：docs/熱磁發電機架構規格_簡報.pptx、docs/熱磁發電機架構規格_詳細.docx
字型 PingFang TC（macOS 中文）。研究數據可信度見內文引用。
"""
from __future__ import annotations

import json
import re
from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import lxml.etree as etree

_SYN_PATHS = [Path("docs/arch_research_synthesis.json"), Path("/tmp/tmg_arch_synth.json")]
_syn_file = next((p for p in _SYN_PATHS if p.exists()), _SYN_PATHS[0])
SYN = json.loads(_syn_file.read_text())
SECTIONS = SYN["sections"]
SPEC = SYN["spec_table"]
CITES = SYN["citations"]
CJK = "PingFang TC"

# GA 架構↔材料共同設計（calibrated 預設；P/η 為整機模型上界，供相對比較）
GA_ROWS = [
    ("近室溫 25°C", "飛輪", "Fe24Al19Co18Cr16", 37, 0.422, 0.55, 1838),
    ("近室溫 25°C", "活塞", "Fe29Cr19Al15Ni15", 35, 0.399, 0.10, 1536),
    ("低溫 150°C", "飛輪", "Fe32Co19Al17Cr14", 158, 0.389, 0.40, 1280),
    ("低溫 150°C", "活塞", "Fe30Co17Ni17Cr15", 159, 0.371, 0.09, 1467),
    ("中溫 350°C", "飛輪", "Fe31Co23Ni17Al14", 361, 0.365, 0.29, 1413),
    ("中溫 350°C", "活塞", "Fe34Co20Ni16Al16", 361, 0.348, 0.08, 1310),
    ("高溫 500°C", "飛輪", "Co47Fe20Al16Cr5", 519, 0.468, 0.36, 1374),
    ("高溫 500°C", "活塞", "Co33Fe27Ni15Al14", 514, 0.367, 0.08, 1362),
]

EXEC = ("本規格比較磁熱發電（TMG）的兩種能量解構：飛輪式（旋轉、空間分離冷熱）與活塞式"
        "（往復／共振、時間交替冷熱）。結論：以工業廢熱（連續、大 ΔT）為目標，飛輪式為量產首選"
        "——其效率優勢約 3.5–5×，幾乎全部來自回熱（ε 0.85–0.95 vs 活塞 <0.3）；只有空間分離才"
        "能做到高 ε。活塞式適合 Phase-0 先導（單元、便宜、共振可放大小 ΔT）。下列參數與數字均附"
        "網路調研引用；GA 共同設計顯示最佳合金會隨架構漂移。")

RECO = ("建議路線：活塞做 Phase-0 先導（去風險、單元、易量測）→ 飛輪做量產（高 ε、平穩類-DC、"
        "可規模化至 kW）。決策規則：要效率＋規模＋連續熱源→飛輪；要最低成本驗證／低品位 ΔT／"
        "可攜→活塞。")


def _lines(body_md: str) -> list[str]:
    out = []
    for ln in body_md.split("\n"):
        ln = ln.strip()
        if ln:
            out.append(ln)
    return out


def _plain(s: str) -> str:
    return s.replace("**", "")


# ───────────────────────── Word ─────────────────────────
def build_docx(path: Path) -> None:
    doc = Document()
    st = doc.styles["Normal"]
    st.font.name = CJK
    st.font.size = Pt(10.5)
    st.element.rPr.rFonts.set(qn("w:eastAsia"), CJK)

    def cjk(run):
        rPr = run._element.get_or_add_rPr()
        rf = rPr.find(qn("w:rFonts"))
        if rf is None:
            rf = OxmlElement("w:rFonts"); rPr.append(rf)
        rf.set(qn("w:eastAsia"), CJK)
        run.font.name = CJK

    def h(t, lvl=1):
        p = doc.add_heading(level=lvl)
        r = p.add_run(t); r.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F); cjk(r)

    def para(t, bold=False):
        p = doc.add_paragraph()
        r = p.add_run(t); r.bold = bold; cjk(r); return p

    def md_para(line: str):
        p = doc.add_paragraph()
        m = re.match(r"\*\*(.+?)\*\*[:：]?\s*(.*)", line)
        if m:
            r = p.add_run(m.group(1) + "："); r.bold = True; cjk(r)
            if m.group(2):
                r2 = p.add_run(_plain(m.group(2))); cjk(r2)
        else:
            r = p.add_run(_plain(line)); cjk(r)

    def table(headers, rows):
        t = doc.add_table(rows=1, cols=len(headers)); t.style = "Light Grid Accent 1"
        for i, x in enumerate(headers):
            c = t.rows[0].cells[i]; c.text = ""
            r = c.paragraphs[0].add_run(x); r.bold = True; cjk(r); r.font.size = Pt(9.5)
        for row in rows:
            cells = t.add_row().cells
            for i, v in enumerate(row):
                cells[i].text = ""
                r = cells[i].paragraphs[0].add_run(str(v)); cjk(r); r.font.size = Pt(9)
        doc.add_paragraph()

    title = doc.add_paragraph(); title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("熱磁發電機架構規格"); r.bold = True; r.font.size = Pt(24)
    r.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F); cjk(r)
    sub = doc.add_paragraph(); sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("飛輪式（旋轉） vs 活塞式（往復） · 網路調研 + GA 共同設計")
    r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x2A, 0x9D, 0x8F); cjk(r)
    doc.add_paragraph()

    h("0. 執行摘要", 1); para(EXEC)
    h("1. 架構↔材料共同設計（GA 結果）", 1)
    para("同一 surrogate，但不同架構的 ε/util/L/B 預設使 GA 選出不同最佳合金與不同 P/η。"
         "（P、η 為整機模型上界，宜相對比較；Tc 為合成模型，配方屬示意。）")
    table(["情境", "架構", "最佳配方 (at%)", "Tc °C", "ΔM (T)", "η %", "P (kW/m³)"],
          [[a, b, c, d, f"{e:.3f}", f"{g:.2f}", f"{hh:.0f}"] for (a, b, c, d, e, g, hh) in GA_ROWS])

    h("2. 飛輪 vs 活塞：整機規格表", 1)
    table(["參數", "飛輪式", "活塞式", "依據"],
          [[s["param"], s["flywheel"], s["piston"], s["basis"]] for s in SPEC])

    h("3. 子系統設計細節", 1)
    for s in SECTIONS:
        h(s["title"], 2)
        for ln in _lines(s["body_md"]):
            md_para(ln)

    h("4. 建議路線", 1); para(RECO, bold=False)

    h("5. 參考文獻（網路調研）", 1)
    for i, c in enumerate(CITES, 1):
        p = doc.add_paragraph(style="List Number")
        r = p.add_run(f"{c.get('title','')} — {c.get('url','')}"); r.font.size = Pt(8.5); cjk(r)

    doc.save(path)
    print("saved", path)


# ───────────────────────── PPT ─────────────────────────
def build_pptx(path: Path) -> None:
    NAVY = PColor(0x1F, 0x3A, 0x5F); TEAL = PColor(0x2A, 0x9D, 0x8F)
    DARK = PColor(0x26, 0x2B, 0x33); GREY = PColor(0x5A, 0x63, 0x70)
    LIGHT = PColor(0xF2, 0xF4, 0xF7); WHITE = PColor(0xFF, 0xFF, 0xFF)
    GOLD = PColor(0xE9, 0xC4, 0x66)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def sf(run, size, bold=False, color=DARK):
        run.font.name = CJK; run.font.size = PPt(size); run.font.bold = bold; run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
        if ea is None:
            ea = etree.SubElement(rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
        ea.set("typeface", CJK)

    def rect(s, x, y, w, h, c):
        sp = s.shapes.add_shape(1, x, y, w, h); sp.fill.solid(); sp.fill.fore_color.rgb = c
        sp.line.fill.background(); sp.shadow.inherit = False; return sp

    def tb(s, x, y, w, h, lines, align=PP_ALIGN.LEFT):
        t = s.shapes.add_textbox(x, y, w, h).text_frame; t.word_wrap = True
        for i, (txt, sz, b, c) in enumerate(lines):
            p = t.paragraphs[0] if i == 0 else t.add_paragraph()
            p.alignment = align; p.space_after = PPt(4); p.line_spacing = 1.08
            r = p.add_run(); r.text = txt; sf(r, sz, b, c)
        return t

    def header(s, kick, title, n):
        rect(s, 0, 0, Inches(0.16), Inches(7.5), TEAL)
        tb(s, Inches(0.5), Inches(0.26), Inches(11.5), Inches(0.4), [(kick, 12, True, TEAL)])
        tb(s, Inches(0.5), Inches(0.54), Inches(12.2), Inches(0.8), [(title, 25, True, NAVY)])
        tb(s, Inches(12.5), Inches(0.26), Inches(0.6), Inches(0.4), [(f"{n:02d}", 12, True, GREY)], PP_ALIGN.RIGHT)

    def slide(bg=WHITE):
        s = prs.slides.add_slide(BLANK)
        if bg != WHITE:
            rect(s, 0, 0, prs.slide_width, prs.slide_height, bg)
        return s

    def bullets(s, x, y, w, h, items, size=15, gap=7):
        t = s.shapes.add_textbox(x, y, w, h).text_frame; t.word_wrap = True
        for i, it in enumerate(items):
            p = t.paragraphs[0] if i == 0 else t.add_paragraph()
            p.space_after = PPt(gap); p.line_spacing = 1.08
            r = p.add_run(); r.text = "•  " + it; sf(r, size, False, DARK)

    def table(s, x, y, col_w, rows, fs=11, rh=0.4):
        yy = y
        for ri, row in enumerate(rows):
            xx = x; h = Inches(rh)
            for ci, cell in enumerate(row):
                w = Inches(col_w[ci])
                bg = NAVY if ri == 0 else (LIGHT if ri % 2 else WHITE)
                rect(s, xx, yy, w, h, bg)
                t = s.shapes.add_textbox(xx + Inches(0.04), yy, w - Inches(0.08), h).text_frame
                t.word_wrap = True; t.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = t.paragraphs[0]; r = p.add_run(); r.text = str(cell)
                sf(r, fs, ri == 0, WHITE if ri == 0 else DARK)
                xx += w
            yy += h

    # cover
    s = slide(NAVY)
    rect(s, 0, Inches(5.0), prs.slide_width, Inches(0.1), GOLD)
    tb(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5), [("THERMOMAGNETIC GENERATOR — ARCHITECTURE SPEC", 14, True, GOLD)])
    tb(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.6),
       [("熱磁發電機架構規格", 42, True, WHITE), ("飛輪式（旋轉） vs 活塞式（往復）", 26, True, LIGHT)])
    tb(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.8),
       [("網路調研（46 引用） + GA 架構↔材料共同設計", 16, False, LIGHT)])
    tb(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4), [("alloy_engine · 2026", 12, False, GOLD)])

    # exec + reco
    s = slide(); header(s, "EXECUTIVE SUMMARY", "摘要與建議", 2)
    tb(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.6), [(EXEC, 16, False, DARK)])
    rect(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(1.7), LIGHT)
    tb(s, Inches(0.8), Inches(4.65), Inches(11.7), Inches(1.5),
       [("建議", 15, True, TEAL), (RECO, 15, False, DARK)])

    # GA results
    s = slide(); header(s, "MATERIAL CO-DESIGN", "架構↔材料共同設計（GA）", 3)
    rows = [["情境", "架構", "最佳配方", "Tc°C", "ΔM", "η%", "P kW/m³"]]
    for (a, b, c, d, e, g, hh) in GA_ROWS:
        rows.append([a, b, c, str(d), f"{e:.2f}", f"{g:.2f}", f"{hh:.0f}"])
    table(s, Inches(0.5), Inches(1.7), [2.0, 1.0, 3.3, 1.2, 1.2, 1.2, 1.6], rows, fs=11, rh=0.46)
    tb(s, Inches(0.5), Inches(6.5), Inches(12.2), Inches(0.7),
       [("飛輪 η 優勢 3.5–5×（源自回熱 ε 0.90 vs 0.20）；最佳合金隨架構漂移。P/η 為模型上界，宜相對比較。", 12, True, GREY)])

    # spec table (split into two slides if long)
    half = (len(SPEC) + 1) // 2
    for part, chunk in enumerate([SPEC[:half], SPEC[half:]]):
        s = slide(); header(s, "SPEC TABLE", f"飛輪 vs 活塞 規格表 ({part+1}/2)", 4 + part)
        rows = [["參數", "飛輪式", "活塞式"]]
        for sp in chunk:
            rows.append([sp["param"], sp["flywheel"][:60], sp["piston"][:60]])
        table(s, Inches(0.5), Inches(1.6), [2.6, 5.1, 5.1], rows, fs=10.5, rh=0.92)

    # 6 subsystem slides
    for i, sec in enumerate(SECTIONS):
        s = slide(); header(s, "SUBSYSTEM", sec["title"], 6 + i)
        items = [_plain(ln)[:120] for ln in _lines(sec["body_md"])][:7]
        bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.4), items, size=14, gap=8)

    # references
    s = slide(); header(s, "REFERENCES", f"參考文獻（{len(CITES)} 篇，節錄）", 12)
    items = [f"{c.get('title','')[:70]} — {c.get('url','')[:48]}" for c in CITES[:14]]
    bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.4), items, size=11, gap=5)

    prs.save(path)
    print("saved", path, "slides", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    Path("docs").mkdir(exist_ok=True)
    build_docx(Path("docs/熱磁發電機架構規格_詳細.docx"))
    build_pptx(Path("docs/熱磁發電機架構規格_簡報.pptx"))
