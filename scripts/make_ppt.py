"""
產生「熱磁發電機材料 → 整機設計」研究簡報（約 30 頁，繁體中文）。
內容全部引自本專案文件與程式碼，無杜撰數字。
輸出：docs/熱磁發電機_研究簡報.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ASSETS = Path("docs/ppt_assets")
OUT = Path("docs/熱磁發電機_研究簡報.pptx")

# 16:9
SW, SH = Inches(13.333), Inches(7.5)

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
ORANGE = RGBColor(0xE7, 0x6F, 0x51)
GOLD = RGBColor(0xE9, 0xC4, 0x66)
DARK = RGBColor(0x26, 0x2B, 0x33)
GREY = RGBColor(0x5A, 0x63, 0x70)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "WenQuanYi Zen Hei"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set_font(run, size, bold=False, color=DARK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 確保東亞字型也套用
    rPr = run._r.get_or_add_rPr()
    import lxml.etree as etree
    ea = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = etree.SubElement(rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface", FONT)


def rect(slide, x, y, w, h, color):
    sp = slide.shapes.add_shape(1, x, y, w, h)  # rectangle
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color, space_after_pt)。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    for i, item in enumerate(lines):
        text, size, bold, color, sa = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        p.line_spacing = 1.12
        # 支援同段多 run（以 \t 分隔粗體前綴）
        run = p.add_run()
        run.text = text
        _set_font(run, size, bold, color)
    return tb


def header(slide, kicker, title, n):
    """頁首：左側色條 + kicker + 標題 + 右上頁碼。"""
    rect(slide, 0, 0, Inches(0.18), SH, TEAL)
    textbox(slide, Inches(0.55), Inches(0.30), Inches(11.5), Inches(0.4),
            [(kicker, 14, True, TEAL, 0)])
    textbox(slide, Inches(0.55), Inches(0.62), Inches(11.8), Inches(0.9),
            [(title, 30, True, NAVY, 0)])
    rect(slide, Inches(0.55), Inches(1.52), Inches(2.2), Pt(3), GOLD)
    # 頁碼
    textbox(slide, Inches(12.4), Inches(0.30), Inches(0.7), Inches(0.4),
            [(f"{n:02d}", 14, True, GREY, 0)], align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, h, items, size=18, gap=10):
    """items: list of (level, text, color?) ; level0=圓點, level1=次項。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        level = it[0]
        text = it[1]
        color = it[2] if len(it) > 2 else DARK
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        bullet = "●  " if level == 0 else "–  "
        run = p.add_run()
        run.text = bullet + text
        _set_font(run, size if level == 0 else size - 2,
                  bold=(level == 0), color=color)
        p.level = level
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    if w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    return slide.shapes.add_picture(str(path), x, y, height=h)


def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    if bg != WHITE:
        rect(s, 0, 0, SW, SH, bg)
    return s


# ───────────────────────── 01 封面 ─────────────────────────
s = new_slide(NAVY)
rect(s, 0, Inches(5.4), SW, Inches(0.12), GOLD)
textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
        [("FIRST-PRINCIPLES DEVICE ENGINEERING", 16, True, GOLD, 0)])
textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.8),
        [("熱磁發電機（TMG）", 48, True, WHITE, 6),
         ("材料篩選 → 整機設計 → 真實資料", 32, True, LIGHT, 0)])
textbox(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(1.0),
        [("從純數學與熱力學底層邏輯出發，建出整機效能模型，", 18, False, LIGHT, 2),
         ("並一路追到真實量測資料校準的研究歷程", 18, False, LIGHT, 0)])
textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
        [("alloy_engine  ·  研究簡報  ·  2026", 14, False, GOLD, 0)])

# ───────────────────────── 02 一句話總覽 ─────────────────────────
s = new_slide()
header(s, "EXECUTIVE SUMMARY", "研究主線（一句話）", 2)
textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.6),
        [("從 TMG 的數學/熱力學推導出發，建出整機效能模型，", 24, False, DARK, 6),
         ("發現瓶頸依序是「顯熱 → 磁滯 → 熱導率 κ → 元素空間 → 真實資料」，", 24, True, NAVY, 6),
         ("並逐一處理；最終把真實 NEMAD 居禮溫度接進 GA 搜尋。", 24, False, DARK, 0)])
add_image(s, ASSETS / "bottleneck.png", Inches(1.4), Inches(4.4), w=Inches(10.5))

# ───────────────────────── 03 由來/緣起 ─────────────────────────
s = new_slide()
header(s, "ORIGIN", "由來：為什麼做這件事", 3)
bullets(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(5),
        [(0, "起點：以「純數學與熱力學底層邏輯」推算並設計熱磁發電機（TMG）。"),
         (1, "把問題拆成三階段：熱能輸入 Q_in、磁矩變化 ΔM、電磁功輸出 W。"),
         (0, "原引擎只是「熱磁合金材料篩選器」——只算材料層。"),
         (1, "材料層回答不了「能組出怎樣的發電機」這個整機問題。"),
         (0, "本研究的工作：把引擎從材料層延伸到整機層，再追到真實資料。"),
         (1, "材料層 → 整機層 → 真實量測校準，三段式升級。"),
         (0, "正反向對偶：發電（TMG）與製冷（磁熱製冷）共用材料與磁路。"),
         (1, "製冷側有成熟文獻（CAS 全固態 HMR）可當量化校準錨點。")])

# ───────────────────────── 04 第一性原理：三階段方程式 ─────────────────────────
s = new_slide()
header(s, "FIRST PRINCIPLES", "整機三階段數學模型", 4)
# 三欄
cols = [
    (TEAL, "① 熱能輸入 Q_in", ["顯熱 + 磁熵兩部分：", "Q = ρCpΔT(1-ε) + ρT·ΔS_M",
                              "ε = 回熱效率", "→ 顯熱主導，回熱是主槓桿"]),
    (ORANGE, "② 磁矩變化 ΔM", ["Maxwell 關係：", "(∂M/∂T)_H = (∂S/∂H)_T",
                              "M(T) 隨溫度跨越 Tc 而變", "→ 需高 ΔM 的工作溫區"]),
    (NAVY, "③ 電磁功輸出 W", ["磁功積分：", "W = μ₀ ∮ H dM",
                             "感應電壓 V_rms ∝ f·N·ΔΦ", "→ 頻率/匝數放大輸出"]),
]
x = Inches(0.7)
cw = Inches(3.95)
for color, title, lines in cols:
    rect(s, x, Inches(2.0), cw, Inches(0.7), color)
    textbox(s, x, Inches(2.05), cw, Inches(0.6), [(title, 18, True, WHITE, 0)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, x, Inches(2.7), cw, Inches(3.2), LIGHT)
    textbox(s, x + Inches(0.2), Inches(2.9), cw - Inches(0.4), Inches(2.9),
            [(l, 15, (i == 0), DARK if i else NAVY, 8) for i, l in enumerate(lines)])
    x = x + cw + Inches(0.25)
textbox(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.8),
        [("效率 η = W / Q_in；卡諾上限 η_C = 1 - T_cold/T_hot。三式單位一致、正反向對偶成立。",
          16, True, GREY, 0)])

# ───────────────────────── 05 架構優化結果 ─────────────────────────
s = new_slide()
header(s, "DEVICE SIMULATION", "整機架構優化：模擬結果", 5)
add_image(s, ASSETS / "architecture.png", Inches(0.7), Inches(1.9), w=Inches(7.4))
bullets(s, Inches(8.4), Inches(2.1), Inches(4.5), Inches(4.5),
        [(0, "分層發電床 + 窄化每層溫區。"),
         (0, "效率 η：0.022% → 1.56%（×71）。"),
         (0, "功率密度：0.73 → 4.1 MW/m³。"),
         (0, "感應電壓：0.14 → 4.37 V（×31）。"),
         (0, "結論：顯熱與磁滯是死穴，回熱是主槓桿。", ORANGE)], size=16, gap=12)

# ───────────────────────── 06 反向：磁熱製冷對偶 ─────────────────────────
s = new_slide()
header(s, "DUALITY", "反向運作：磁熱製冷對偶", 6)
bullets(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4),
        [(0, "同一套材料與磁路，正向發電、反向製冷——熱力學對偶。"),
         (1, "製冷 COP = Q_cold / W；發電 η = W / Q_in。"),
         (0, "對標 CAS 全固態磁熱製冷 HMR：8.4 kW/kg、火用效率 59.6%。"),
         (1, "與文獻報告 8.3 kW/kg、54.2% 量級吻合 → 模型可信。"),
         (0, "關鍵發現：磁滯（hysteresis）是發電與製冷兩個方向的頭號殺手。", ORANGE),
         (1, "一階材料 ΔM 大但磁滯損耗也大 → 需 hysteresis-penalized 評分。")])

# ───────────────────────── 07 隱形殺手 κ ─────────────────────────
s = new_slide()
header(s, "HIDDEN BOTTLENECK", "材料 what-if：熱導率 κ 是隱形殺手", 7)
bullets(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.5),
        [(0, "直覺：只要找 ΔM 最大的材料就好——錯。"),
         (0, "整機循環頻率 f ≈ α/(2L²)，α = κ/(ρCp) 是熱擴散率。"),
         (1, "κ 太低 → 熱來不及進出材料 → 循環頻率被卡死 → 功率上不去。"),
         (0, "最好的磁熱材料（一階）往往 κ 很低，形成隱形天花板。", ORANGE),
         (0, "因此整機級目標不能只看 ΔM，必須同時看 ΔM × κ。", NAVY)])

# ───────────────────────── 08 複合材料 ─────────────────────────
s = new_slide()
header(s, "COMPOSITE", "解法：複合材料（高 κ 基底 + 高 ΔM 相）", 8)
bullets(s, Inches(0.9), Inches(1.9), Inches(6.0), Inches(4.5),
        [(0, "把高 ΔM 的磁熱相嵌入高 κ 基底（Cu / Al / α-Fe）。"),
         (0, "Wiener 上下界給出有效 κ 的可行範圍。"),
         (0, "對每個候選找最佳基底分率 φ*。"),
         (0, "功率密度提升 ×2.4–×10（D4 封頂後的物理量級）。", TEAL),
         (1, "封頂前曾膨脹到 ×22，屬非物理假象（見 D4）。")], size=16)
rect(s, Inches(7.3), Inches(2.1), Inches(5.4), Inches(3.6), LIGHT)
textbox(s, Inches(7.5), Inches(2.3), Inches(5.0), Inches(3.2),
        [("複合有效物性", 18, True, NAVY, 10),
         ("κ_eff：介於並聯/串聯界之間", 15, False, DARK, 8),
         ("ΔM_eff = (1-φ)·ΔM_相", 15, False, DARK, 8),
         ("φ* = 使 功率密度 最大的基底分率", 15, False, DARK, 8),
         ("→ 瓶頸從「材料 κ」轉移到「元素空間」", 15, True, ORANGE, 0)])

# ───────────────────────── 09 GA 接整機目標 ─────────────────────────
s = new_slide()
header(s, "OPTIMIZATION", "GA 直接最佳化整機目標", 9)
bullets(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.5),
        [(0, "GPU 向量化遺傳演算法，族群可達數萬，每秒數百萬次評估。"),
         (0, "兩種 fitness 模式：softmag（軟磁）與 thermomagnetic（熱磁）。"),
         (0, "整機目標 w_device > 0 時，GA 直接最佳化「功率密度 × 效率」。", NAVY),
         (1, "--device-matrix 讓 GA 為每個候選自動找最佳複合基底分率。"),
         (0, "化學可合成性軟約束：脆相、析出相、鐵磁基底不足、稀土可製造性。"),
         (1, "讓 GA 不選出物理上無法製造的配方。")])

# ───────────────────────── 10 稀土擴張 ─────────────────────────
s = new_slide()
header(s, "SEARCH SPACE", "元素空間擴張：加入稀土 Gd / La", 10)
bullets(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.5),
        [(0, "瓶頸轉移到元素空間後，擴張 10 → 12 元素（加 Gd、La）。"),
         (1, "Gd：低 Tc 鐵磁體（≈293K / 20°C），室溫磁熱效應的標竿材料。"),
         (1, "La：搭配 Fe-Si 形成 La(Fe,Si)₁₃ 一階磁熱相，Tc 近室溫。"),
         (0, "新增「近室溫廢熱_25C」情境，對應稀土室溫 MCE。"),
         (0, "結果：GA 在該情境自主收斂到 Gd 基室溫材料——行為符合物理。", TEAL)])

# ───────────────────────── 11 缺陷評估 ─────────────────────────
s = new_slide()
header(s, "HONEST ASSESSMENT", "持平評估：12 項已知缺陷", 11)
textbox(s, Inches(0.9), Inches(1.85), Inches(11.8), Inches(0.7),
        [("整體判斷：概念完整、內部一致、適合相對比較與設計探索；但絕對量化可信度被一個根因卡住——", 16, False, DARK, 0)])
rect(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(0.8), ORANGE)
textbox(s, Inches(1.1), Inches(2.68), Inches(11.1), Inches(0.7),
        [("根因：以「合成物理」訓練，未對標真實量測。", 20, True, WHITE, 0)],
        anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(3),
        [(0, "第一級（可信度）：D1 Tc 落差、D2 稀土外推、D3 ΔM 無真實來源。"),
         (0, "第二級（物理保真度）：D4 頻率無上限、D5 一階相變、D6 假設參數、D7 無 CI。"),
         (0, "第三級（覆蓋/工程）：D8 缺 P/Ge、D9 可製造性高估、D10–D12。")], size=15, gap=8)

# ───────────────────────── 12 缺陷狀態總覽圖 ─────────────────────────
s = new_slide()
header(s, "REMEDIATION", "缺陷處理狀態：5 項已修", 12)
add_image(s, ASSETS / "defect_status.png", Inches(1.5), Inches(1.9), w=Inches(10.3))

# ───────────────────────── 13 D7 CI ─────────────────────────
s = new_slide()
header(s, "FIX · D7", "D7：建立 CI 迴歸防護", 13)
bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4),
        [(0, "問題：10→12 元素擴張曾「靜默」弄壞 NEMAD 腳本（維度不符），無人察覺。"),
         (0, "修復：GitHub Actions 在每次 push / PR 跑全套件 + 7 支腳本 py_compile。"),
         (0, "效果：把沉默迴歸變成 CI 紅燈，第一時間攔截。", TEAL),
         (0, "這是「預防」類缺陷的根本防護，故列為最高優先處理。", NAVY)])

# ───────────────────────── 14 D4 頻率封頂 ─────────────────────────
s = new_slide()
header(s, "FIX · D4", "D4：整機頻率封頂（κ 不再永遠有益）", 14)
bullets(s, Inches(0.9), Inches(2.0), Inches(6.1), Inches(4),
        [(0, "問題：f = α/(2L²) 無上限 → κ 越高 f 越高 → 複合增益膨脹到 ×22（非物理）。"),
         (0, "修復：有效頻率封頂", NAVY),
         (1, "f_eff = f / (1 + f / f_max)，預設 f_max = 50 Hz。"),
         (0, "效果：複合增益收斂到 ×2.4–×10 的物理量級。", TEAL)], size=16)
rect(s, Inches(7.3), Inches(2.2), Inches(5.4), Inches(2.2), LIGHT)
textbox(s, Inches(7.5), Inches(2.5), Inches(5.0), Inches(1.8),
        [("f_eff = f / (1 + f / f_max)", 22, True, NAVY, 12),
         ("f → 0 時 f_eff ≈ f", 16, False, DARK, 6),
         ("f → ∞ 時 f_eff → f_max（飽和）", 16, False, DARK, 0)])

# ───────────────────────── 15 D5 一階相變 ─────────────────────────
s = new_slide()
header(s, "FIX · D5", "D5：一階相變的銳變（logistic）", 15)
add_image(s, ASSETS / "mt_curves.png", Inches(0.7), Inches(1.9), w=Inches(7.2))
bullets(s, Inches(8.1), Inches(2.2), Inches(4.8), Inches(4),
        [(0, "平均場 M=Ms(1-T/Tc)^0.5 低估一階材料。"),
         (0, "改用 logistic 階梯：", NAVY),
         (1, "M = Ms / (1+exp((T-Tc)/w))。"),
         (0, "transition_width_K 參數控制銳度。"),
         (1, "預設 None = 二階，向後相容。"),
         (0, "一階 ΔM > 平均場，描述更正確。", TEAL)], size=15, gap=8)

# ───────────────────────── 16 真實資料：重大發現 ─────────────────────────
s = new_slide()
header(s, "REAL DATA", "重大發現：真實資料其實「拿得到」", 16)
bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4),
        [(0, "NEMAD FM_with_curie.csv 就在公開 GitHub repo（sumanitani/NEMAD-MagneticML）。"),
         (1, "15,577 筆鐵磁化合物 + 居禮溫度，含 La/Ce/Pr/Nd/Gd 稀土，欄位與腳本完全吻合。"),
         (0, "Materials Project API：DFT 磁化（需免費 API key），已驗證可連線。"),
         (0, "授權：NEMAD repo 無 license → 採「執行時抓取 + 引用」，資料不入版控。", ORANGE),
         (0, "→ 把缺陷 D1 從「被阻擋」改判為「現在就能做」。", TEAL)])

# ───────────────────────── 17 sim-to-real 證據 ─────────────────────────
s = new_slide()
header(s, "PROOF", "硬證據：合成資料在真實資料上失效", 17)
add_image(s, ASSETS / "sim_to_real.png", Inches(1.3), Inches(1.95), w=Inches(10.6))
textbox(s, Inches(0.9), Inches(6.5), Inches(11.8), Inches(0.7),
        [("合成代理 R² = -0.17（比直接猜平均還差）；用真實 NEMAD 訓練 R² = 0.88，與論文報告的 0.86 一致。",
          16, True, NAVY, 0)])

# ───────────────────────── 18 稀土外推崩潰 ─────────────────────────
s = new_slide()
header(s, "PROOF · D2", "稀土外推的硬證據", 18)
bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4),
        [(0, "Fe-Co 系尚可：Hiperco50 誤差 -19°C、Permalloy +30°C。"),
         (0, "稀土/特殊系崩潰：", ORANGE),
         (1, "La-Fe-Si 合成預測 +463°C，NEMAD 真值 -57°C → 誤差 +520°C。"),
         (1, "Invar Fe65Ni35 誤差 +420°C。"),
         (0, "意義：合成物理無法外推到稀土；但 NEMAD 內就有正確的 La-Fe-Si Tc（≈216K）。", NAVY),
         (1, "→ 真實資料同時是問題的證據，也是解答的來源。", TEAL)])

# ───────────────────────── 19 MP 磁化 / Br 偏差 ─────────────────────────
s = new_slide()
header(s, "PROOF · D3", "MP 磁化：量出 Br 系統性低估", 19)
bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4),
        [(0, "用 MP DFT 磁化對標 8 個已知鐵磁體的合成 Br。"),
         (1, "MAE = 0.53 T；系統性偏差 bias = -0.50 T（合成 Br 系統性低估）。"),
         (0, "意義：本引擎的 ΔM 與整機功率密度是系統性偏保守（低估）。", NAVY),
         (0, "謹慎處理：MP 是 0K DFT 飽和，合成 Br 是工作溫度——", ORANGE),
         (1, "直接拿 MP 值「校正」會把 ΔM 推高到不對的方向 → 列為謹慎後續，不做天真校正。")])

# ───────────────────────── 20 HybridBundle ─────────────────────────
s = new_slide()
header(s, "INTEGRATION", "把真實 Tc 接進 GA：HybridBundle", 20)
bullets(s, Inches(0.9), Inches(2.0), Inches(6.0), Inches(4),
        [(0, "真實 NEMAD 只有 Tc，沒有 Hc/Br/σy。"),
         (0, "HybridBundle = 真實 Tc 模型 + 合成 Hc/Br/σy。", NAVY),
         (0, "對外介面與 SurrogateBundle 完全相同。"),
         (1, "GA 無需改動即可用真實居禮溫度搜尋。"),
         (0, "--hybrid-tc 旗標即可啟用。", TEAL)], size=16)
rect(s, Inches(7.3), Inches(2.2), Inches(5.4), Inches(3.4), LIGHT)
textbox(s, Inches(7.5), Inches(2.45), Inches(5.0), Inches(3.0),
        [("predict_properties(comp)：", 16, True, NAVY, 10),
         ("  Tc       ← 真實 NEMAD 模型", 15, False, TEAL, 8),
         ("  Hc       ← 合成代理", 15, False, DARK, 8),
         ("  Br       ← 合成代理", 15, False, DARK, 8),
         ("  strength ← 合成代理", 15, False, DARK, 10),
         ("→ 真實 Tc 命中 153–157°C 目標帶", 15, True, NAVY, 0)])

# ───────────────────────── 21 D9 稀土可製造性（最新） ─────────────────────────
s = new_slide()
header(s, "FIX · D9 (NEW)", "D9：稀土可製造性懲罰（本次新增）", 21)
add_image(s, ASSETS / "rare_earth_penalty.png", Inches(0.7), Inches(1.9), w=Inches(7.0))
bullets(s, Inches(8.0), Inches(2.1), Inches(4.9), Inches(4.5),
        [(0, "問題：GA 化學約束不罰稀土，Gd 還被當鐵磁基底獎勵。"),
         (0, "稀土易氧化、一階相脆裂，可製造性被高估。"),
         (0, "修復：兩個漸進懲罰", NAVY),
         (1, "氧化：(Gd+La) 線性，純稀土→0.75。"),
         (1, "脆裂：稀土×(Fe+Si) 交互。"),
         (0, "down-rank 但不排除有效 MCE。", TEAL)], size=14, gap=7)

# ───────────────────────── 22 扎實 vs 限制 ─────────────────────────
s = new_slide()
header(s, "WHAT'S SOLID", "哪些扎實、哪些是限制（持平記錄）", 22)
rect(s, Inches(0.7), Inches(2.0), Inches(5.9), Inches(0.65), TEAL)
textbox(s, Inches(0.7), Inches(2.05), Inches(5.9), Inches(0.6),
        [("✅ 扎實", 18, True, WHITE, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.8), Inches(2.8), Inches(5.7), Inches(3.8),
        [(0, "整機熱力學/電磁學方程正確、單位一致。"),
         (0, "正反向（發電↔製冷）對偶成立。"),
         (0, "製冷側對標文獻 HMR 量級吻合。"),
         (0, "設計結論被 GA 自主行為佐證。"),
         (0, "真實 Tc R²=0.88，可信度已實證。")], size=15, gap=9)
rect(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(0.65), ORANGE)
textbox(s, Inches(6.8), Inches(2.05), Inches(5.9), Inches(0.6),
        [("⚠ 限制", 18, True, WHITE, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(6.9), Inches(2.8), Inches(5.7), Inches(3.8),
        [(0, "Hc/Br/σy 仍為合成（待 MP 校準）。"),
         (0, "Br 系統性低估 0.50 T。"),
         (0, "ΔM 銳度 w 仍需真實資料校準。"),
         (0, "複合微結構參數為假設值。"),
         (0, "發電側尚無真實原型對標。")], size=15, gap=9)

# ───────────────────────── 23 系統架構 ─────────────────────────
s = new_slide()
header(s, "ARCHITECTURE", "系統架構：從元素到整機", 23)
boxes = [
    (TEAL, "元素 / 成分", "12 元素 Dirichlet 取樣"),
    (NAVY, "代理模型", "Tc / Hc / Br / σy（合成 + 真實 Tc）"),
    (ORANGE, "整機 / 複合", "功率密度 × 效率，複合 κ"),
    (GOLD, "GA 搜尋", "GPU 向量化，化學約束"),
]
y = Inches(2.2)
for color, t, d in boxes:
    rect(s, Inches(1.2), y, Inches(3.2), Inches(0.9), color)
    textbox(s, Inches(1.2), y, Inches(3.2), Inches(0.9), [(t, 18, True, WHITE, 0)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.7), y, Inches(7.8), Inches(0.9), [(d, 17, False, DARK, 0)],
            anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(1.1)
textbox(s, Inches(1.2), Inches(6.7), Inches(11), Inches(0.6),
        [("CI（GitHub Actions）跨層守護：222 項測試 + 腳本編譯，每次 push 驗證。", 15, True, GREY, 0)])

# ───────────────────────── 24 真實資料來源 ─────────────────────────
s = new_slide()
header(s, "DATA SOURCES", "真實資料來源與授權", 24)
bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5),
        [(0, "NEMAD FM_with_curie.csv：15,577 FM 化合物 + Curie 溫度（含稀土）。"),
         (1, "公開 GitHub 可抓；引用 Nature Communications 2025；放 external/（git-ignored）。"),
         (0, "Materials Project：DFT 磁化 / 磁序。"),
         (1, "需免費 API key；存 external/.mp_key（git-ignored）或 GitHub Secret。"),
         (0, "精選 2,504 筆 Curie 資料集（arXiv 2509.17464）可作對照。"),
         (0, "原則：資料不入版控，只「執行時抓取 + 引用」，繞過授權疑慮。", NAVY)])

# ───────────────────────── 25 GHA 自動化方案 ─────────────────────────
s = new_slide()
header(s, "AUTOMATION", "GitHub Actions 可自動化方案", 25)
bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5),
        [(0, "方案 A（最高槓桿，現在就能做）：排程抓公開 NEMAD + 跑 sim-to-real 對標 → 存 artifact。", TEAL),
         (1, "讓「合成 vs 真實」落差變成持續可見的指標，且不需 commit 資料。"),
         (0, "方案 B：連續基準（benchmark-action），追蹤 R²/MAE 隨 commit 變化，回歸即標記。"),
         (0, "方案 C：資料完整性——SHA256 digest + pandas 欄位/範圍驗證。"),
         (0, "方案 D：MP 機密——API key 存 GitHub Secret，安全拉磁化資料。")])

# ───────────────────────── 26 復現指令 ─────────────────────────
s = new_slide(LIGHT)
header(s, "REPRODUCE", "復現指令", 26)
cmds = [
    "# 1. 合成代理（含 Gd/La 稀土）",
    "python scripts/train_surrogate.py --n-samples 6000 --epochs 100",
    "",
    "# 2. 真實 NEMAD Tc 基準  → R² ≈ 0.88",
    "python scripts/train_surrogate_nemad_baseline.py",
    "",
    "# 3. sim-to-real 對標",
    "python scripts/nemad_eval.py",
    "python scripts/mp_magnetization_eval.py   # 需 MP key",
    "",
    "# 4. GA 用真實 Tc + 整機/複合目標搜尋",
    "python scripts/run_search.py --scenario 低溫廢熱_150C \\",
    "    --mode thermomagnetic --hybrid-tc <nemad_baseline.pt> \\",
    "    --w-device 1.0 --device-matrix Cu",
]
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.6), Inches(5))
tf = tb.text_frame
tf.word_wrap = True
for i, c in enumerate(cmds):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(4)
    run = p.add_run()
    run.text = c if c else " "
    run.font.name = "DejaVu Sans Mono"
    run.font.size = Pt(15)
    run.font.bold = c.startswith("#")
    run.font.color.rgb = TEAL if c.startswith("#") else DARK

# ───────────────────────── 27 量化結果總表 ─────────────────────────
s = new_slide()
header(s, "KEY NUMBERS", "關鍵量化結果", 27)
rows = [
    ("主題", "數字", True),
    ("架構優化（低溫情境）", "η 0.022%→1.56%（×71）、電壓 ×31", False),
    ("複合材料（Mn-Fe-P+Cu）", "功率密度 ×10（D4 封頂後物理量級）", False),
    ("製冷對標（CAS HMR）", "8.4 kW/kg、火用效率 59.6%", False),
    ("Tc sim-to-real", "合成 R²=-0.17 → 真實 R²=0.88（MAE 274→81°C）", False),
    ("Br vs MP DFT", "bias -0.50 T（系統性低估）", False),
    ("測試 / CI", "222 項測試通過，GitHub runner 綠燈", False),
]
y = Inches(2.0)
for label, val, hdr in rows:
    bg = NAVY if hdr else (LIGHT if (rows.index((label, val, hdr)) % 2) else WHITE)
    rect(s, Inches(0.8), y, Inches(11.7), Inches(0.62), bg)
    c1 = WHITE if hdr else DARK
    textbox(s, Inches(1.0), y, Inches(4.6), Inches(0.62),
            [(label, 16, hdr, c1, 0)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(5.7), y, Inches(6.6), Inches(0.62),
            [(val, 16, hdr, c1, 0)], anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.66)

# ───────────────────────── 28 現況 ─────────────────────────
s = new_slide()
header(s, "STATUS", "現況", 28)
bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5),
        [(0, "全部 PR 已合併進 main（整機模型 → 真實 Tc → 知識備份 → D9）。", TEAL),
         (0, "已完成：整機模型、製冷對偶、複合材料、稀土擴張、CI、D4/D5 物理修正、"
             "真實 Tc 接入 GA、D9 稀土可製造性。"),
         (0, "進行中：真實 Br 校準（需區分 0K 飽和 vs 工作溫度磁化）。", ORANGE),
         (0, "已連通：NEMAD（Tc）、MP（磁化）皆可取用，剩整合與校準。"),
         (0, "222 項測試通過，CI 於 GitHub runner 實證綠燈。", NAVY)])

# ───────────────────────── 29 下一步 ─────────────────────────
s = new_slide()
header(s, "NEXT", "下一步路線圖", 29)
items = [
    (TEAL, "近期", "D2：把真實 NEMAD Tc 併入主代理（不只 inference 期的 HybridBundle）。"),
    (NAVY, "近期", "GHA 方案 A：排程 NEMAD sim-to-real 基準 → artifact。"),
    (GOLD, "中期", "D3：建 0K-飽和 → 工作溫度的 Br 溫度修正模型，再併入 MP 磁化。"),
    (ORANGE, "中期", "D8：補 P/Ge 元素與 La-Fe-Si-H 氫化 Tc 上修模型。"),
    (GREY, "長期", "D12：找文獻 TMG 原型數據，錨定發電側 W/η/P。"),
]
y = Inches(2.0)
for color, tag, txt in items:
    rect(s, Inches(0.9), y, Inches(1.4), Inches(0.8), color)
    textbox(s, Inches(0.9), y, Inches(1.4), Inches(0.8), [(tag, 16, True, WHITE, 0)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(2.5), y, Inches(10), Inches(0.8), [(txt, 16, False, DARK, 0)],
            anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.95)

# ───────────────────────── 30 結語 ─────────────────────────
s = new_slide(NAVY)
rect(s, 0, Inches(3.5), SW, Inches(0.1), GOLD)
textbox(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.6),
        [("從第一性原理到真實資料", 40, True, WHITE, 8),
         ("一條完整、誠實、可復現的研究線", 26, False, LIGHT, 0)])
bullets_box = s.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11.5), Inches(2.8))
tf = bullets_box.text_frame
tf.word_wrap = True
for i, t in enumerate([
    "整機模型讓我們看清瓶頸鏈：顯熱 → 磁滯 → κ → 元素空間 → 真實資料。",
    "每一步的結論都被物理、文獻或 GA 自主行為佐證；缺陷誠實登錄。",
    "真實 NEMAD Tc 把可信度從「相對工具」推向「絕對預測」（R² 0.88）。",
    "下一步：把真實磁化謹慎校準進來，並用 GHA 持續監測 sim-to-real 落差。",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    run = p.add_run()
    run.text = "●  " + t
    _set_font(run, 17, False, LIGHT)

prs.save(str(OUT))
print("saved", OUT, "with", len(prs.slides._sldIdLst), "slides")
